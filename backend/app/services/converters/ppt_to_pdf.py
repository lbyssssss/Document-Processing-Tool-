# PPT to PDF Converter
from pathlib import Path
from typing import Optional


class PPTToPDFConverter:
    """PPT转PDF转换器"""

    def __init__(self, ppt_path: Path, output_path: Path):
        self.ppt_path = ppt_path
        self.output_path = output_path

    def convert(self) -> dict:
        """将PowerPoint演示文稿转换为PDF

        注意：此实现提供基础框架。完整的PPT到PDF转换通常需要：
        1. Windows + Microsoft Office (使用win32com)
        2. LibreOffice (使用subprocess调用)
        3. 云服务API

        当前实现使用reportlab创建基础PDF，保留文本内容
        """
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.pdfgen import canvas
            from reportlab.lib import colors
            from reportlab.lib.units import inch
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            # 加载PPT文件
            prs = Presentation(str(self.ppt_path))

            # 创建PDF
            c = canvas.Canvas(str(self.output_path), pagesize=A4)

            # 页面尺寸
            page_width, page_height = A4

            # 为每一张幻灯片创建一个PDF页面
            for slide_num, slide in enumerate(prs.slides):
                # 开始新页面（除了第一页）
                if slide_num > 0:
                    c.showPage()

                # 绘制幻灯片标题
                slide_title = self._get_slide_title(slide)
                if slide_title:
                    c.setFont("Helvetica-Bold", 16)
                    c.drawString(50, page_height - 50, slide_title)

                # 绘制幻灯片内容
                y_position = page_height - 80
                line_height = 20

                # 获取幻灯片中的所有文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text:
                        text = shape.text.strip()
                        if text and text != slide_title:
                            # 检查是否需要新页面
                            if y_position < 50:
                                c.showPage()
                                y_position = page_height - 50

                            c.setFont("Helvetica", 10)
                            c.drawString(50, y_position, text)
                            y_position -= line_height

            # 保存PDF
            c.save()

            return {
                "success": True,
                "output_file": str(self.output_path),
                "slides_converted": len(prs.slides),
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        try:
            # 尝试从标题占位符获取
            for shape in slide.shapes:
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if hasattr(shape, "placeholder_format"):
                        pf = shape.placeholder_format
                        if pf.type == 1:  # 标题
                            if hasattr(shape, "text_frame"):
                                return shape.text_frame.text.strip()

            # 如果没有找到标题占位符，返回第一个文本框
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text:
                    return shape.text_frame.text.strip()[:100]

            return ""
        except:
            return ""

    def convert_with_libreoffice(self) -> dict:
        """使用LibreOffice转换PPT到PDF

        需要系统安装LibreOffice
        """
        try:
            import subprocess
            from tempfile import TemporaryDirectory

            with TemporaryDirectory() as tmpdir:
                # 使用LibreOffice转换
                cmd = [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmpdir,
                    str(self.ppt_path)
                ]

                result = subprocess.run(cmd, capture_output=True, text=True)

                if result.returncode != 0:
                    return {
                        "success": False,
                        "error": f"LibreOffice转换失败: {result.stderr}",
                    }

                # 查找生成的PDF文件
                import shutil
                pdf_files = list(Path(tmpdir).glob("*.pdf"))

                if not pdf_files:
                    return {
                        "success": False,
                        "error": "未找到生成的PDF文件",
                    }

                # 移动PDF到目标位置
                shutil.move(str(pdf_files[0]), str(self.output_path))

                return {
                    "success": True,
                    "output_file": str(self.output_path),
                    "method": "libreoffice",
                }

        except FileNotFoundError:
            return {
                "success": False,
                "error": "LibreOffice未安装",
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    def convert_with_win32com(self) -> dict:
        """使用Windows COM接口转换PPT到PDF

        需要Windows系统和Microsoft Office
        """
        try:
            import win32com.client

            # 创建PowerPoint应用
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False

            try:
                # 打开PPT文件
                presentation = powerpoint.Presentations.Open(str(self.ppt_path.absolute()))

                # 导出为PDF
                presentation.SaveAs(str(self.output_path.absolute()), 32)  # 32 = ppSaveAsPDF

                # 关闭演示文稿
                presentation.Close()

                return {
                    "success": True,
                    "output_file": str(self.output_path),
                    "method": "win32com",
                }

            finally:
                # 退出PowerPoint应用
                powerpoint.Quit()

        except ImportError:
            return {
                "success": False,
                "error": "pywin32未安装",
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    def get_slide_count(self) -> int:
        """获取PPT幻灯片数量"""
        try:
            from pptx import Presentation
            prs = Presentation(str(self.ppt_path))
            return len(prs.slides)
        except:
            return 0
