# PDF to PPT Converter
from pathlib import Path
from typing import Optional
from PyPDF2 import PdfReader


class PDFToPPTConverter:
    """PDF转PPT转换器"""

    def __init__(self, pdf_path: Path, output_path: Path):
        self.pdf_path = pdf_path
        self.output_path = output_path

    def convert(self) -> dict:
        """将PDF转换为PowerPoint演示文稿

        注意：此实现提供基础框架，完整的PDF到PPT转换需要更复杂的布局分析
        """
        try:
            reader = PdfReader(str(self.pdf_path))

            from pptx import Presentation
            from pptx.util import Inches, Pt

            # 创建演示文稿
            prs = Presentation()

            # 获取幻灯片尺寸（默认16:9）
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 为每一页PDF创建一个幻灯片
            for page_num, page in enumerate(reader.pages):
                # 使用空白布局
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)

                # 提取页面文本
                text_content = self._extract_page_text(page)

                # 添加文本框
                if text_content:
                    self._add_text_to_slide(slide, text_content, slide_width, slide_height)

                # 可以在此添加PDF页面作为图片（需要pdf2image库）
                # self._add_pdf_page_as_image(slide, page, page_num)

            # 保存PPT文件
            prs.save(str(self.output_path))

            return {
                "success": True,
                "output_file": str(self.output_path),
                "slides_created": len(reader.pages),
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    def _extract_page_text(self, page) -> str:
        """提取PDF页面的文本内容"""
        try:
            text = page.extract_text()
            return text if text else ""
        except:
            return ""

    def _add_text_to_slide(self, slide, text: str, slide_width, slide_height):
        """向幻灯片添加文本框"""
        from pptx.util import Inches, Pt

        # 创建文本框（居中，留边距）
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        # 设置边距
        text_frame.word_wrap = True

        # 分割文本为段落
        lines = text.split('\n')

        for line in lines:
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(12)
                p.font.name = "宋体"

    def _add_pdf_page_as_image(self, slide, page, page_num):
        """将PDF页面作为图片添加到幻灯片

        注意：需要pdf2image库
        """
        try:
            from pdf2image import convert_from_path
            import io

            # 将PDF页面转换为图片
            images = convert_from_path(
                str(self.pdf_path),
                first_page=page_num + 1,
                last_page=page_num + 1,
                dpi=150
            )

            if images:
                img = images[0]
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='PNG')
                img_byte_arr.seek(0)

                # 添加图片到幻灯片
                slide.shapes.add_picture(
                    img_byte_arr,
                    Inches(0),
                    Inches(0),
                    width=Inches(10)
                )
        except ImportError:
            pass
        except Exception as e:
            print(f"添加图片失败: {e}")

    def get_page_count(self) -> int:
        """获取PDF页面数量"""
        try:
            reader = PdfReader(str(self.pdf_path))
            return len(reader.pages)
        except:
            return 0
