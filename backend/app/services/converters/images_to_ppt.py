# Images to PPT Conversion
from PIL import Image
from pathlib import Path
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import re


class ImagesToPPTConverter:
    """图片转PPT转换器"""

    def __init__(self, images: List[Path], output_path: Path, ocr_mode: bool = False):
        self.images = images
        self.output_path = output_path
        self.ocr_mode = ocr_mode

    def convert(self) -> dict:
        """将多张图片转换为PPT文档"""
        try:
            if not self.images:
                return {
                    "success": False,
                    "error": "没有提供图片文件",
                }

            # 创建PPT演示文稿
            prs = Presentation()

            # 处理每张图片
            images_added = []

            for i, image_path in enumerate(self.images):
                result = self._add_image_to_ppt(prs, image_path, i)
                if result.get("success"):
                    images_added.append(result)
                else:
                    return {
                        "success": False,
                        "error": f"添加图片{image_path.name}失败: {result.get('error')}",
                    }

            # 保存PPT文件
            prs.save(str(self.output_path))

            return {
                "success": True,
                "output_file": str(self.output_path),
                "images_added": len(images_added),
                "images": images_added,
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    def _add_image_to_ppt(self, prs, image_path: Path, index: int) -> dict:
        """添加单张图片到PPT"""
        try:
            # 打开图片
            img = Image.open(image_path)

            # 获取图片尺寸
            img_width, img_height = img.size

            # 如果是OCR模式，提取文字并添加到PPT
            if self.ocr_mode:
                return self._add_ocr_text_to_ppt(prs, image_path, img_width, img_height, index)
            else:
                return self._add_image_to_slide(prs, image_path, img_width, img_height, index)

        except Exception as e:
            return {
                "success": False,
                "error": f"添加图片{image_path.name}失败: {str(e)}",
            }

    def _add_ocr_text_to_ppt(self, prs, image_path: Path, img_width: int, img_height: int, index: int) -> dict:
        """OCR模式：识别图片中的文字并按格式添加到PPT"""
        try:
            from pytesseract import image_to_string, image_to_data
            from pytesseract import Output

            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 使用 Tesseract OCR 识别文字
            ocr_data = image_to_data(str(image_path), lang='chi_sim+eng', output_type=Output.DICT)

            # 获取幻灯片尺寸
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 解析识别到的文本，提取格式信息
            if ocr_data and 'text' in ocr_data:
                text_list = ocr_data['text']
                if text_list:
                    # 分析文本块，提取格式和内容
                    parsed_blocks = self._analyze_text_blocks(text_list, ocr_data)

                    # 按原始位置排序文本块
                    parsed_blocks.sort(key=lambda x: x["original_top"])

                    # 按格式添加到PPT
                    for block in parsed_blocks:
                        self._add_text_block_to_slide(
                            slide,
                            block,
                            img_width,
                            img_height,
                            slide_width
                        )

                    return {
                        "success": True,
                        "image_name": image_path.name,
                        "mode": "OCR",
                        "text_lines": len(parsed_blocks),
                    }

            # 未识别到文字，添加占位文本
            text_box = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(4)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = f"未识别到文字内容\n\n图片: {image_path.name}"
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"

            return {
                "success": True,
                "image_name": image_path.name,
                "mode": "OCR",
                "text_lines": 0,
            }

        except ImportError:
            return {
                "success": False,
                "error": "OCR模块未安装，请运行: apt-get install tesseract-ocr && pip install pytesseract",
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"OCR识别失败: {str(e)}",
            }

    def _add_text_block_to_slide(self, slide, block: dict, img_width: int, img_height: int, slide_width: int, slide_height: int):
        """添加文本块到幻灯片"""
        text = block["text"]

        # 计算位置（映射到PPT坐标）
        ppt_left = (block["left"] / img_width) * slide_width
        ppt_top = (block["top"] / img_height) * slide_height
        ppt_width = (block["width"] / img_width) * slide_width

        # 根据文本类型设置字体
        if block["type"] == "title":
            font_size = Pt(32)
            font_name = "微软雅黑"
            font_bold = True
            font_color = RGBColor(0, 51, 102)
        elif block["type"] == "list":
            font_size = Pt(18)
            font_name = "微软雅黑"
            font_bold = False
            font_color = RGBColor(50, 50, 50)
        elif block["type"] == "content":
            font_size = Pt(14)
            font_name = "微软雅黑"
            font_bold = False
            font_color = RGBColor(0, 0, 0)
        else:
            font_size = Pt(14)
            font_name = "微软雅黑"
            font_bold = False
            font_color = RGBColor(0, 0, 0)

        # 添加文本框
        text_box = slide.shapes.add_textbox(
            left=ppt_left,
            top=ppt_top,
            width=max(ppt_width, Inches(2)),
            height=max(Pt(24), (block["height"] / img_height) * slide_height),
        )

        # 设置文本格式
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.name = font_name
        p.font.bold = font_bold
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.LEFT

    def _analyze_text_blocks(self, text_list: list, ocr_data: dict) -> list:
        """分析文本块，提取格式信息"""
        blocks = []
        i = 0

        while i < len(text_list):
            text = text_list[i]
            if not text.strip():
                i += 1
                continue

            # 获取位置信息
            if i < len(ocr_data.get('left', [])):
                left = ocr_data['left'][i]
                top = ocr_data['top'][i]
                width = ocr_data['width'][i]
                height = ocr_data['height'][i]
            else:
                left = top = width = height = 0

            # 分析文本特征判断类型
            text_type, block_content = self._classify_text_block(text, i)

            if text_type == "empty":
                i += 1
                continue

            blocks.append({
                "text": text,
                "type": text_type,
                "original_top": top,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "content": block_content,
            })

            i += 1

        return blocks

    def _classify_text_block(self, text: str, index: int) -> Tuple[str, str]:
        """分类文本块类型"""
        text = text.strip()

        # 空行
        if not text:
            return "empty", ""

        # 判断是否是标题
        # 1. 全大写或首字母大写
        # 2. 短文本（通常<20字符）
        # 3. 行首有编号或符号
        if text.isupper():
            return "title", text

        if len(text) < 30:
            # 检查是否是编号标题
            if re.match(r'^[一二三四五六七八九十]+[、，\.]', text):
                return "title", text
            # 检查是否以符号开头
            if re.match(r'^[•\-\*]+', text):
                return "title", text

        # 判断是否是列表项
        # 1. 有编号或符号开头
        if re.match(r'^[一二三四五六七八九十]+[、\.\)]+', text):
            return "list", text

        # 正常内容
        return "content", text

    def _add_image_to_slide(self, prs, image_path: Path, img_width: int, img_height: int, index: int) -> dict:
        """普通模式：直接添加图片到PPT"""
        try:
            # 计算适合PPT页面的尺寸（使用16:9比例）
            ppt_width, ppt_height = self._calculate_fit_dimensions(img_width, img_height)

            # 将图片转换为PPT可用格式
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 添加图片到幻灯片
            slide.shapes.add_picture(
                img_buffer,
                left=Inches(0),
                top=Inches(0),
                width=Inches(ppt_width),
                height=Inches(ppt_height)
            )

            return {
                "success": True,
                "image_name": image_path.name,
                "mode": "IMAGE",
                "original_size": f"{img_width}x{img_height}",
                "fitted_size": f"{ppt_width}x{ppt_height}",
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"添加图片{image_path.name}失败: {str(e)}",
            }

    def _calculate_fit_dimensions(self, img_width: int, img_height: int) -> tuple[float, float]:
        """计算图片适合PPT幻灯片的尺寸"""
        # PPT标准幻灯片比例 16:9
        ppt_aspect = 16 / 9
        img_aspect = img_width / img_height

        # 使用固定尺寸（10英寸宽）
        base_width = 10

        if img_aspect > ppt_aspect:
            # 图片更宽，以宽度为基准
            ppt_width = base_width
            ppt_height = base_width / img_aspect
        else:
            # 图片更高，以高度为基准
            ppt_height = base_width / ppt_aspect
            ppt_width = ppt_height * img_aspect

        return (ppt_width, ppt_height)

    def get_images_info(self) -> dict:
        """获取图片信息"""
        try:
            images_info = []
            for image_path in self.images:
                with Image.open(image_path) as img:
                    images_info.append({
                        "name": image_path.name,
                        "width": img.size[0],
                        "height": img.size[1],
                        "format": img.format,
                        "mode": img.mode,
                    })
            return {
                "success": True,
                "images": images_info,
                "total_images": len(images_info),
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }
